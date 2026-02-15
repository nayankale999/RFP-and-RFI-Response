import io
import logging

from pptx import Presentation

from app.documents.parsers.base import BaseParser, ParsedDocument

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Microsoft PowerPoint (.pptx) parser."""

    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    def parse(self, file_data: bytes, filename: str) -> ParsedDocument:
        try:
            prs = Presentation(io.BytesIO(file_data))
        except Exception as e:
            logger.error(f"Failed to parse PPTX {filename}: {e}")
            raise

        slides_text = []
        all_tables = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    all_tables.append(rows)

            slides_text.append("\n".join(slide_parts))

        full_text = "\n\n".join(slides_text)
        return ParsedDocument(
            text=full_text,
            page_count=len(prs.slides),
            tables=all_tables,
            metadata={
                "filename": filename,
                "parser": "pptx",
                "slide_count": len(prs.slides),
            },
        )
